from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


ROOT = Path(r"C:\Users\dell\Desktop\HA\HA")
OUT = ROOT / "output" / "poster_ha"
PNG = OUT / "HA_algorithm_poster_v3.png"
PPTX = OUT / "HA_algorithm_poster_v3.pptx"


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(20)
    prs.slide_height = Inches(50)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(PNG), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(str(PPTX))
    print(PPTX)


if __name__ == "__main__":
    main()
